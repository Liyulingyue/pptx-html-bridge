import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .converters import emu_to_px, color_to_hex, dash_style_to_css

def collect_layout_elements(layout, layout_index_map, html_dir, slide_width_px, slide_height_px, background_style):
    """Collect layout images and shapes, and update background_style if full-slide image."""
    layout_images = []
    layout_shapes = []
    try:
        layout_idx = layout_index_map.get(id(layout), 0)
        # Extract layout picture shapes and simple layout shapes (LINE/AUTO_SHAPE)
        for li, lshape in enumerate(layout.shapes):
            if lshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    limg = lshape.image
                    lbytes = limg.blob
                    lext = limg.ext
                    lfname = f"layout{layout_idx}_img{li}.{lext}"
                    lpath = os.path.join(html_dir, lfname)
                    if not os.path.exists(lpath):
                        with open(lpath, 'wb') as _f:
                            _f.write(lbytes)
                    lleft = emu_to_px(lshape.left)
                    ltop = emu_to_px(lshape.top)
                    lw = emu_to_px(lshape.width)
                    lh = emu_to_px(lshape.height)
                    layout_images.append((lfname, lleft, ltop, lw, lh))
                except Exception:
                    pass
            elif lshape.shape_type == MSO_SHAPE_TYPE.LINE or lshape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Extract style for line/auto shape on layout
                try:
                    lleft = emu_to_px(lshape.left)
                    ltop = emu_to_px(lshape.top)
                    lw = emu_to_px(lshape.width)
                    lh = emu_to_px(lshape.height)
                    rot = getattr(lshape, 'rotation', 0) or 0
                    shape_info = {'type': lshape.shape_type, 'left': lleft, 'top': ltop, 'width': lw, 'height': lh, 'rotation': rot}
                    # line stroke
                    try:
                        if hasattr(lshape, 'line') and lshape.line is not None:
                            lw_px = 1
                            if lshape.line.width:
                                # line.width returns EMU; convert to px using emu_to_px
                                lw_px_calc = emu_to_px(lshape.line.width)
                                lw_px = max(1, lw_px_calc)
                            shape_info['stroke_width'] = lw_px
                            if hasattr(lshape.line, 'color') and hasattr(lshape.line.color, 'rgb'):
                                shape_info['stroke_color'] = color_to_hex(lshape.line.color)
                            # dash style
                            try:
                                ds = getattr(lshape.line, 'dash_style', None)
                                shape_info['dash_style'] = dash_style_to_css(ds)
                            except Exception:
                                pass
                    except Exception:
                        pass
                    # fill color for auto shapes
                    if lshape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        try:
                            if lshape.fill and hasattr(lshape.fill, 'fore_color') and hasattr(lshape.fill.fore_color, 'rgb'):
                                shape_info['fill_color'] = color_to_hex(lshape.fill.fore_color)
                        except Exception:
                            pass
                    layout_shapes.append(shape_info)
                except Exception:
                    pass
        # Extract master picture shapes and master shapes
        try:
            master = layout.slide_master
            for mi, mshape in enumerate(master.shapes):
                if mshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        mimg = mshape.image
                        mbytes = mimg.blob
                        mext = mimg.ext
                        mfname = f"master_img{mi}.{mext}"
                        mpath = os.path.join(html_dir, mfname)
                        if not os.path.exists(mpath):
                            with open(mpath, 'wb') as _f:
                                _f.write(mbytes)
                        mleft = emu_to_px(mshape.left)
                        mtop = emu_to_px(mshape.top)
                        mw = emu_to_px(mshape.width)
                        mh = emu_to_px(mshape.height)
                        layout_images.append((mfname, mleft, mtop, mw, mh))
                    except Exception:
                        pass
                elif mshape.shape_type == MSO_SHAPE_TYPE.LINE or mshape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    try:
                        mleft = emu_to_px(mshape.left)
                        mtop = emu_to_px(mshape.top)
                        mw = emu_to_px(mshape.width)
                        mh = emu_to_px(mshape.height)
                        mrot = getattr(mshape, 'rotation', 0) or 0
                        mshape_info = {'type': mshape.shape_type, 'left': mleft, 'top': mtop, 'width': mw, 'height': mh, 'rotation': mrot}
                        try:
                            if hasattr(mshape, 'line') and mshape.line is not None:
                                mw_px = 1
                                if mshape.line.width:
                                    mw_px_calc = emu_to_px(mshape.line.width)
                                    mw_px = max(1, mw_px_calc)
                                mshape_info['stroke_width'] = mw_px
                                if hasattr(mshape.line, 'color') and hasattr(mshape.line.color, 'rgb'):
                                    mshape_info['stroke_color'] = color_to_hex(mshape.line.color)
                                try:
                                    mshape_info['dash_style'] = dash_style_to_css(getattr(mshape.line, 'dash_style', None))
                                except Exception:
                                    pass
                        except Exception:
                            pass
                        if mshape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            try:
                                if mshape.fill and hasattr(mshape.fill, 'fore_color') and hasattr(mshape.fill.fore_color, 'rgb'):
                                    mshape_info['fill_color'] = color_to_hex(mshape.fill.fore_color)
                            except Exception:
                                pass
                        layout_shapes.append(mshape_info)
                    except Exception:
                        pass
        except Exception:
            pass
    except Exception:
        layout_images = []

    # If a layout image covers the full slide, use it as slide background
    layout_images_filtered = []
    for lfname, lleft, ltop, lw, lh in layout_images:
        if lleft <= 1 and ltop <= 1 and lw >= slide_width_px - 2 and lh >= slide_height_px - 2:
            # treat as background
            background_style = f"background-image: url('{lfname}'); background-size: cover; background-repeat: no-repeat; background-position: center;"
        else:
            layout_images_filtered.append((lfname, lleft, ltop, lw, lh))

    return layout_images_filtered, layout_shapes, background_style