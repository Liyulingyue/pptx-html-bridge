import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from .converters import emu_to_px, emu_to_pt, color_to_hex, pt_to_px
from .fonts import get_effective_font

def html_builder():
    """Create a simple html line builder for pretty printing."""
    lines = []
    def add(line, indent=0):
        lines.append(('    ' * indent) + line)
    def to_str(compact=False):
        if compact:
            return ''.join(line.lstrip() for line in lines)
        return '\n'.join(lines)
    return add, to_str

def generate_index_html(filename_base, num_slides, html_dir, compact):
    """Generate index.html for the presentation slides."""
    add_idx, to_str_idx = html_builder()
    add_idx('<!DOCTYPE html>')
    add_idx('<html lang="zh-CN">')
    add_idx('<head>', 1)
    add_idx('<meta charset="UTF-8">', 2)
    add_idx('<title>PPT Slides Index</title>', 2)
    add_idx('</head>', 1)
    add_idx('<body>', 1)
    add_idx('<h1>幻灯片列表</h1>', 2)
    add_idx('<ul>', 2)
    for i in range(1, num_slides + 1):
        add_idx(f'<li><a href="slides/slide{i}.html">Slide {i}</a></li>', 3)
    add_idx('</ul>', 2)
    add_idx('</body>', 1)
    add_idx('</html>')
    with open(os.path.join(html_dir, f"{filename_base}_index.html"), 'w', encoding='utf-8') as f:
        f.write(to_str_idx(compact=compact))

def generate_main_html(source_dir, html_dir, compact):
    """Generate main.html entry page."""
    main_add, main_to_str = html_builder()
    main_add('<!DOCTYPE html>')
    main_add('<html lang="zh-CN">')
    main_add('<head>', 1)
    main_add('<meta charset="UTF-8">', 2)
    main_add('<title>PPT to HTML Converter - Main Entry</title>', 2)
    main_add('<style>', 2)
    main_add('body { font-family: Arial, sans-serif; margin: 40px; background-color: #f4f4f4; }', 3)
    main_add('.container { max-width: 800px; margin: 0 auto; background: white; padding: 20px; border-radius: 8px; box-shadow: 0 0 10px rgba(0,0,0,0.1); }', 3)
    main_add('h1 { color: #333; }', 3)
    main_add('ul { list-style-type: none; padding: 0; }', 3)
    main_add('li { margin: 10px 0; }', 3)
    main_add('a { text-decoration: none; color: #007bff; font-size: 18px; }', 3)
    main_add('a:hover { text-decoration: underline; }', 3)
    main_add('</style>', 2)
    main_add('</head>', 1)
    main_add('<body>', 1)
    main_add('<div class="container">', 2)
    main_add('<h1>PPT to HTML 转换结果</h1>', 3)
    main_add('<p>以下是转换后的演示文稿：</p>', 3)
    main_add('<ul>', 3)
    for filename in os.listdir(source_dir):
        if filename.endswith('.pptx'):
            filename_base = os.path.splitext(filename)[0]
            main_add(f'<li><a href="{filename_base}_index.html">{filename_base}</a></li>', 4)
    main_add('</ul>', 3)
    main_add('<p>点击链接查看幻灯片。</p>', 3)
    main_add('</div>', 2)
    main_add('</body>', 1)
    main_add('</html>')
    with open(os.path.join(html_dir, 'main.html'), 'w', encoding='utf-8') as f:
        f.write(main_to_str(compact=compact))

def generate_slide_html(i, num_slides, theme_minor_font, slide_width_px, slide_height_px, background_style, nav, layout_images_filtered, layout_shapes, slide, prs, layout_placeholder_defaults, html_dir, compact):
    """Generate HTML for a single slide."""
    add, to_str = html_builder()
    # Prepare fallback font-family: theme minor font -> Chinese fallback -> Arial -> sans-serif
    default_font_stack = []
    if theme_minor_font:
        default_font_stack.append(theme_minor_font)
    # Add common Chinese fonts and system fonts for better fidelity
    default_font_stack.extend(["微软雅黑", "Microsoft YaHei", "Helvetica", "Arial", "sans-serif"])
    default_font_family = ', '.join([f'"{f}"' for f in default_font_stack])
    # top part
    add('<!DOCTYPE html>')
    add('<html lang="zh-CN">')
    add('<head>', 1)
    add('<meta charset="UTF-8">', 2)
    add(f'<title>Slide {i}</title>', 2)
    add('<style>', 2)
    add(f'body {{ font-family: {default_font_family}; padding: 20px; }}', 3)
    add(f'.slide {{ position: relative; width: {slide_width_px}px; height: {slide_height_px}px; {background_style} border: 1px solid #ccc; margin: 0 auto; box-sizing: border-box; overflow: hidden; }}', 3)
    add('.shape { position: absolute; z-index: 2; box-sizing: border-box; }', 3)
    add('.layout-image { position: absolute; z-index: 0; }', 3)
    add('.layout-shape { position: absolute; z-index: 1; box-sizing: border-box; }', 3)
    add('.shape img { display: block; object-fit: contain; }', 3)
    add('* { -webkit-font-smoothing: antialiased; text-rendering: optimizeLegibility; }', 3)
    add('p { line-height: 1.15; margin: 0; }', 3)
    add('table { border-collapse: collapse; }', 3)
    add('td, th { border: 1px solid #000; padding: 4px; }', 3)
    add('.nav { text-align: center; margin-bottom: 20px; }', 3)
    add('</style>', 2)
    add('</head>', 1)
    add('<body>', 1)
    add(nav, 2)
    add('<div class="slide">', 2)
    add('<!-- layout/master images -->', 3)
    # Append layout images (non-full-slide)
    for lfname, lleft, ltop, lw, lh in layout_images_filtered:
        lstyle = f"left: {lleft}px; top: {ltop}px; width: {lw}px; height: {lh}px;"
        add(f'<div class="shape layout-image" style="{lstyle}"><img src="{lfname}" style="width: 100%; height: 100%;" alt="Background Image"></div>', 3)
    # render layout shapes (lines / auto shapes)
    for lshape in layout_shapes:
        try:
            stype = lshape.get('type')
            sleft = lshape.get('left')
            stop = lshape.get('top')
            sw = lshape.get('width')
            sh = lshape.get('height')
            srot = lshape.get('rotation', 0)
            if stype == MSO_SHAPE_TYPE.LINE:
                stroke_width = lshape.get('stroke_width', 2)
                stroke_color = lshape.get('stroke_color', '#000') or '#000'
                dash_style = lshape.get('dash_style') if lshape.get('dash_style', None) else 'solid'
                sstyle = f"left: {sleft}px; top: {stop}px; width: {sw}px; height: {max(1, stroke_width)}px; transform-origin: left top; transform: rotate({srot}deg);"
                # use border-top for dashed style; if dashed, set border-top style else fill
                css_border = f"border-top: {stroke_width}px {dash_style} {stroke_color};"
                add(f'<div class="shape layout-shape line" style="{sstyle} {css_border}"></div>', 3)
            elif stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
                fill_color = lshape.get('fill_color', 'transparent') or 'transparent'
                stroke_color = lshape.get('stroke_color')
                stroke_width = lshape.get('stroke_width')
                border_style = ''
                if stroke_color and stroke_width:
                    border_style = f"border: {stroke_width}px solid {stroke_color};"
                sstyle = f"left: {sleft}px; top: {stop}px; width: {sw}px; height: {sh}px; background-color: {fill_color}; {border_style}; transform-origin: left top; transform: rotate({srot}deg);"
                add(f'<div class="shape layout-shape auto-shape" style="{sstyle}"></div>', 3)
        except Exception:
            pass
    
    img_count = 0
    for shape in slide.shapes:
        left_px = emu_to_px(shape.left)
        top_px = emu_to_px(shape.top)
        width_px = emu_to_px(shape.width)
        height_px = emu_to_px(shape.height)
        shape_style = f"left: {left_px}px; top: {top_px}px; width: {width_px}px; height: {height_px}px;"
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Handle images
            image = shape.image
            image_bytes = image.blob
            ext = image.ext
            img_filename = f"slide{i}_img{img_count}.{ext}"
            # Save image to media directory (passed as html_dir parameter)
            media_dir = os.path.join(os.path.dirname(html_dir), "media")
            with open(os.path.join(media_dir, img_filename), 'wb') as f:
                f.write(image_bytes)
            add(f'<div class="shape" style="{shape_style}"><img src="../media/{img_filename}" style="width: 100%; height: 100%;" alt="Image"></div>', 3)
            img_count += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Handle tables
            table = shape.table
            table_html = "<table>"
            for row in table.rows:
                table_html += "<tr>"
                for cell in row.cells:
                    table_html += f"<td>{cell.text.replace(chr(13), '<br>')}</td>"
                table_html += "</tr>"
            table_html += "</table>"
            add(f'<div class="shape" style="{shape_style}">{table_html}</div>', 3)
        elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            # Handle videos
            try:
                # Find video relationship
                video_rel = None
                for rel in shape.part.rels.values():
                    if 'video' in rel.reltype or 'media' in rel.reltype:
                        video_rel = rel
                        break
                
                if video_rel and video_rel.target_part and hasattr(video_rel.target_part, 'blob'):
                    # Extract video file
                    video_bytes = video_rel.target_part.blob
                    # Determine file extension from content type or filename
                    ext = 'mp4'  # default
                    if hasattr(video_rel.target_part, 'content_type'):
                        if 'mp4' in video_rel.target_part.content_type:
                            ext = 'mp4'
                        elif 'avi' in video_rel.target_part.content_type:
                            ext = 'avi'
                        elif 'mov' in video_rel.target_part.content_type:
                            ext = 'mov'
                        elif 'wmv' in video_rel.target_part.content_type:
                            ext = 'wmv'
                    
                    video_filename = f"slide{i}_video{img_count}.{ext}"
                    # Save video to media directory
                    media_dir = os.path.join(os.path.dirname(html_dir), "media")
                    with open(os.path.join(media_dir, video_filename), 'wb') as f:
                        f.write(video_bytes)
                    
                    # Generate video HTML with poster frame if available
                    poster_attr = ""
                    if shape.poster_frame:
                        poster_bytes = shape.poster_frame.blob
                        poster_ext = shape.poster_frame.ext
                        poster_filename = f"slide{i}_poster{img_count}.{poster_ext}"
                        with open(os.path.join(media_dir, poster_filename), 'wb') as f:
                            f.write(poster_bytes)
                        poster_attr = f' poster="../media/{poster_filename}"'
                    
                    video_html = f'<video controls style="width: 100%; height: 100%;"{poster_attr}><source src="../media/{video_filename}" type="video/{ext}">Your browser does not support the video tag.</video>'
                    add(f'<div class="shape" style="{shape_style}">{video_html}</div>', 3)
                    img_count += 1
            except Exception as e:
                # Fallback: just show a placeholder
                add(f'<div class="shape" style="{shape_style}"><div style="width: 100%; height: 100%; background: #f0f0f0; display: flex; align-items: center; justify-content: center; border: 1px solid #ccc;">[Video]</div></div>', 3)
        elif hasattr(shape, "text_frame") and shape.text_frame:
            # Handle text shapes with full styling
            text_html = ""
            try:
                for paragraph in shape.text_frame.paragraphs:
                    para_style = ""
                    # Get paragraph level properties
                    if paragraph.alignment:
                        align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                        para_style += f"text-align: {align_map.get(paragraph.alignment, 'left')}; "
                    # indentation for bullet/levels
                    try:
                        level = getattr(paragraph, 'level', 0) or 0
                        if level and level > 0:
                            para_style += f"margin-left: {level * 28}px; "
                    except Exception:
                        pass
                    
                    para_html = f'<p style="{para_style}">'
                    for run in paragraph.runs:
                        run_style = ""
                        # detect title placeholder heuristics for default size
                        try:
                            default_is_title = False
                            if getattr(shape, 'is_placeholder', False):
                                ph = shape.placeholder
                                if ph is not None and getattr(ph, 'placeholder_format', None) is not None:
                                    ptype = getattr(ph.placeholder_format, 'type', None)
                                    if ptype in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                                        default_is_title = True
                            # also use spatial heuristic: top near top and text short
                            if not default_is_title:
                                if top_px is not None and top_px < (slide_height_px * 0.18) and len(run.text.strip()) < 25:
                                    default_is_title = True
                        except Exception:
                            default_is_title = False
                        # get effective font family and size based on layout fallback
                        layout_defaults = None
                        try:
                            lid = id(shape.slide.slide_layout)
                            layout_defaults = layout_placeholder_defaults.get(lid, {}).get(getattr(shape.placeholder_format, 'type', None), None)
                        except Exception:
                            layout_defaults = None
                        ff, fsize = get_effective_font(run, paragraph, shape, theme_minor_font, layout_default=layout_defaults, default_is_title=default_is_title)
                        if fsize:
                            fsize_px = pt_to_px(fsize) or int(fsize)
                            run_style += f"font-size: {fsize_px}px; "
                            # approximate line-height based on font size
                            try:
                                line_h_px = int(float(fsize) * 1.15)
                                run_style += f"line-height: {line_h_px}px; "
                            except Exception:
                                pass
                        if ff:
                            run_style += f"font-family: {ff}; "
                        if run.font.bold:
                            run_style += "font-weight: bold; "
                        elif layout_defaults and layout_defaults.get('bold'):
                            run_style += "font-weight: bold; "
                        if run.font.italic:
                            run_style += "font-style: italic; "
                        elif layout_defaults and layout_defaults.get('italic'):
                            run_style += "font-style: italic; "
                        if run.font.underline:
                            run_style += "text-decoration: underline; "
                        elif layout_defaults and layout_defaults.get('underline'):
                            run_style += "text-decoration: underline; "
                        # Extract text color with better fallback logic
                        text_color = None
                        
                        # First try: direct RGB color
                        if run.font.color and hasattr(run.font.color, 'rgb'):
                            try:
                                rgb = run.font.color.rgb
                                if rgb:  # Only if rgb is not None/empty
                                    text_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                            except:
                                pass
                        
                        # Second try: theme/scheme color
                        if text_color is None and run.font.color and hasattr(run.font.color, 'theme_color'):
                            try:
                                theme_color = run.font.color.theme_color
                                if theme_color is not None:
                                    # Try to get RGB from theme color
                                    try:
                                        rgb = run.font.color.rgb
                                        if rgb:
                                            text_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                    except:
                                        pass
                            except Exception:
                                pass
                        
                        # Third try: layout defaults
                        if text_color is None and layout_defaults and layout_defaults.get('color'):
                            text_color = layout_defaults.get('color')
                        
                        # Fourth try: check for solid fill color in run properties
                        if text_color is None:
                            try:
                                # Check if run has color information through other means
                                if hasattr(run.font.color, '_color') and run.font.color._color:
                                    color_obj = run.font.color._color
                                    if hasattr(color_obj, 'rgb') and color_obj.rgb:
                                        rgb = color_obj.rgb
                                        text_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                            except Exception:
                                pass
                        
                        # Apply the color if found
                        if text_color:
                            run_style += f"color: {text_color}; "
                        else:
                            # Last resort: default to white for visibility on dark backgrounds
                            run_style += "color: #ffffff; "
                        if run.font.name:
                            run_style += f"font-family: {run.font.name}; "
                        para_html += f'<span style="{run_style}">{run.text}</span>'
                    para_html += '</p>'
                    text_html += para_html
            except Exception as e:
                # Fallback: just use the text
                text_html = f'<p>{shape.text}</p>'
            
            add(f'<div class="shape" style="{shape_style}">{text_html}</div>', 3)
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            # draw a simple line as a thin rectangle with stroke color
            try:
                stroke_color = '#000'
                stroke_width = 2
                if hasattr(shape, 'line') and shape.line is not None:
                    if hasattr(shape.line, 'color') and hasattr(shape.line.color, 'rgb'):
                        stroke_color = color_to_hex(shape.line.color)
                    if shape.line.width:
                        wpt = emu_to_pt(shape.line.width)
                        if wpt:
                            stroke_width = max(1, int(wpt / 1.333))
                rot = getattr(shape, 'rotation', 0) or 0
                sstyle = f"left: {left_px}px; top: {top_px}px; width: {width_px}px; height: {max(1, stroke_width)}px; background-color: {stroke_color}; transform-origin: left top; transform: rotate({rot}deg);"
                add(f'<div class="shape line" style="{sstyle}"></div>', 3)
            except Exception:
                pass
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                fill_color = 'transparent'
                stroke_color = None
                stroke_width = None
                if shape.fill and hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                    fill_color = color_to_hex(shape.fill.fore_color)
                if hasattr(shape, 'line') and shape.line is not None:
                    if hasattr(shape.line, 'color') and hasattr(shape.line.color, 'rgb'):
                        stroke_color = color_to_hex(shape.line.color)
                    if shape.line.width:
                        wpt = emu_to_pt(shape.line.width)
                        if wpt:
                            stroke_width = max(1, int(wpt / 1.333))
                border_style = ''
                if stroke_color and stroke_width:
                    border_style = f"border: {stroke_width}px solid {stroke_color};"
                rot = getattr(shape, 'rotation', 0) or 0
                sstyle = f"left: {left_px}px; top: {top_px}px; width: {width_px}px; height: {height_px}px; background-color: {fill_color}; {border_style}; transform-origin: left top; transform: rotate({rot}deg);"
                add(f'<div class="shape auto-shape" style="{sstyle}"></div>', 3)
            except Exception:
                pass
    
    add('</div>', 2)
    add('</body>', 1)
    add('</html>')

    slide_filename = os.path.join(html_dir, f"slide{i}.html")
    with open(slide_filename, 'w', encoding='utf-8') as f:
        f.write(to_str(compact=compact))