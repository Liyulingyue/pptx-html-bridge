"""
PPTX to HTML Converter

This module provides functionality to convert PowerPoint (.pptx) files to HTML format.
"""

import os
import sys
from typing import Optional, Dict, Any

from pptx import Presentation
from .utils import *


class PPTXToHTMLConverter:
    """Main converter class for PPTX to HTML conversion."""

    def __init__(self, source_dir: str = None, html_dir: str = None, compact: bool = False):
        """
        Initialize the converter.

        Args:
            source_dir: Directory containing PPTX files (optional)
            html_dir: Directory to output HTML files (optional)
            compact: Whether to generate compact HTML without line breaks
        """
        self.source_dir = source_dir
        self.html_dir = html_dir
        self.compact = compact

    def convert_file(self, pptx_path: str, output_dir: Optional[str] = None) -> Dict[str, Any]:
        """
        Convert a single PPTX file to HTML.

        Args:
            pptx_path: Path to the PPTX file
            output_dir: Output directory (overrides self.html_dir if provided)

        Returns:
            Dict containing conversion results and metadata
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

        html_dir = output_dir or self.html_dir
        if not html_dir:
            html_dir = os.path.splitext(pptx_path)[0] + "_html"

        os.makedirs(html_dir, exist_ok=True)

        # Load the presentation
        prs = Presentation(pptx_path)

        # Determine theme fonts
        try:
            theme_major_font, theme_minor_font = get_theme_fonts(prs)
        except Exception:
            theme_major_font, theme_minor_font = None, None

        filename_base = os.path.splitext(os.path.basename(pptx_path))[0]
        num_slides = len(prs.slides)
        slide_width_px = emu_to_px(prs.slide_width)
        slide_height_px = emu_to_px(prs.slide_height)

        # Map layouts to index for unique naming
        layout_index_map = {id(layout): idx for idx, layout in enumerate(prs.slide_layouts, start=1)}

        # Extract layout placeholder defaults per layout
        layout_placeholder_defaults = {}
        for layout in prs.slide_layouts:
            try:
                layout_placeholder_defaults[id(layout)] = get_layout_placeholder_defaults(layout, prs)
            except Exception:
                layout_placeholder_defaults[id(layout)] = {}

        # Create structured output directories
        slides_dir = os.path.join(html_dir, "slides")
        media_dir = os.path.join(html_dir, "media")
        os.makedirs(slides_dir, exist_ok=True)
        os.makedirs(media_dir, exist_ok=True)

        # Generate index.html
        generate_index_html(filename_base, num_slides, html_dir, self.compact)

        # Process each slide
        generated_files = []
        for i, slide in enumerate(prs.slides, 1):
            # Get background style
            bg_info = get_background_style(slide, prs)
            if bg_info.startswith("picture:"):
                # Handle picture background
                parts = bg_info.split(':')
                bg_img_bytes = eval(parts[1].split(',')[0])
                bg_ext = parts[1].split(',')[1]
                bg_filename = f"slide{i}_bg.{bg_ext}"
                with open(os.path.join(media_dir, bg_filename), 'wb') as f:
                    f.write(bg_img_bytes)
                background_style = f"background-image: url('media/{bg_filename}'); background-size: cover; background-repeat: no-repeat; background-position: center;"
                generated_files.append(f"media/{bg_filename}")
            else:
                background_style = bg_info

            # Collect layout images and master images
            layout_images_filtered, layout_shapes, background_style = collect_layout_elements(
                slide.slide_layout, layout_index_map, media_dir, slide_width_px, slide_height_px, background_style
            )

            # Update layout image paths to include media/ prefix
            layout_images_filtered = [(f"../media/{fname}", left, top, w, h) for fname, left, top, w, h in layout_images_filtered]

            # Update background style to use media/ prefix if it contains a picture
            if "url('" in background_style and not background_style.startswith("background-image: url('media/"):
                # Replace url('filename') with url('media/filename')
                import re
                background_style = re.sub(r"url\('([^']+)'\)", r"url('../media/\1')", background_style)

            # Create navigation (relative paths within slides directory)
            prev_link = f'<a href="slide{i-1}.html">上一页</a>' if i > 1 else ''
            next_link = f'<a href="slide{i+1}.html">下一页</a>' if i < num_slides else ''
            nav = f'<div class="nav">{prev_link} {next_link}</div>'

            # Generate slide HTML in slides directory
            generate_slide_html(
                i, num_slides, theme_minor_font, slide_width_px, slide_height_px,
                background_style, nav, layout_images_filtered, layout_shapes,
                slide, prs, layout_placeholder_defaults, slides_dir, self.compact
            )
            generated_files.append(f"slides/slide{i}.html")

        # Generate index file in root
        index_file = f"{filename_base}_index.html"
        generated_files.append(index_file)

        return {
            "pptx_file": pptx_path,
            "output_dir": html_dir,
            "slides_count": num_slides,
            "generated_files": generated_files,
            "index_file": os.path.join(html_dir, index_file)
        }

    def convert_directory(self, source_dir: Optional[str] = None, output_dir: Optional[str] = None) -> Dict[str, Any]:
        """
        Convert all PPTX files in a directory to HTML.

        Args:
            source_dir: Source directory (overrides self.source_dir if provided)
            output_dir: Output directory (overrides self.html_dir if provided)

        Returns:
            Dict containing conversion results for all files
        """
        src_dir = source_dir or self.source_dir
        if not src_dir:
            raise ValueError("Source directory must be specified")

        if not os.path.exists(src_dir):
            raise FileNotFoundError(f"Source directory not found: {src_dir}")

        html_dir = output_dir or self.html_dir
        if not html_dir:
            html_dir = os.path.join(src_dir, "html_output")

        os.makedirs(html_dir, exist_ok=True)

        results = []
        pptx_files = [f for f in os.listdir(src_dir) if f.endswith('.pptx')]

        for filename in pptx_files:
            pptx_path = os.path.join(src_dir, filename)
            try:
                result = self.convert_file(pptx_path, html_dir)
                results.append(result)
                print(f"Converted {filename} to HTML")
            except Exception as e:
                print(f"Failed to convert {filename}: {e}")
                results.append({
                    "pptx_file": pptx_path,
                    "error": str(e)
                })

        # Generate main index if multiple files
        if len(results) > 1:
            generate_main_html(src_dir, html_dir, self.compact)
            main_index = os.path.join(html_dir, "main.html")
            print(f"Generated main index: {main_index}")

        return {
            "source_dir": src_dir,
            "output_dir": html_dir,
            "converted_files": len([r for r in results if "error" not in r]),
            "failed_files": len([r for r in results if "error" in r]),
            "results": results
        }


def convert_pptx_to_html(pptx_path: str, output_dir: Optional[str] = None, compact: bool = False) -> Dict[str, Any]:
    """
    Convenience function to convert a single PPTX file to HTML.

    Args:
        pptx_path: Path to the PPTX file
        output_dir: Output directory (optional)
        compact: Whether to generate compact HTML

    Returns:
        Dict containing conversion results
    """
    converter = PPTXToHTMLConverter(compact=compact)
    return converter.convert_file(pptx_path, output_dir)


def convert_pptx_directory(source_dir: str, output_dir: Optional[str] = None, compact: bool = False) -> Dict[str, Any]:
    """
    Convenience function to convert all PPTX files in a directory.

    Args:
        source_dir: Source directory containing PPTX files
        output_dir: Output directory (optional)
        compact: Whether to generate compact HTML

    Returns:
        Dict containing conversion results
    """
    converter = PPTXToHTMLConverter(compact=compact)
    return converter.convert_directory(source_dir, output_dir)


def main():
    """Command line interface."""
    import argparse

    parser = argparse.ArgumentParser(description='Convert PPTX files to HTML (one slide per HTML).')
    parser.add_argument('input', help='Input PPTX file or directory')
    parser.add_argument('--output', '-o', help='Output directory')
    parser.add_argument('--compact', action='store_true', help='Write compact HTML (no line breaks, useful for minimal output).')

    args = parser.parse_args()

    converter = PPTXToHTMLConverter(compact=args.compact)

    input_path = args.input
    if os.path.isfile(input_path) and input_path.endswith('.pptx'):
        # Convert single file
        result = converter.convert_file(input_path, args.output)
        print(f"Converted {os.path.basename(input_path)} to HTML")
        print(f"Output directory: {result['output_dir']}")
        print(f"Generated {len(result['generated_files'])} files")
    elif os.path.isdir(input_path):
        # Convert directory
        result = converter.convert_directory(input_path, args.output)
        print(f"Converted {result['converted_files']} files from {input_path}")
        if result['failed_files'] > 0:
            print(f"Failed to convert {result['failed_files']} files")
        print(f"Output directory: {result['output_dir']}")
    else:
        print(f"Error: {input_path} is not a valid PPTX file or directory")
        return 1

    return 0


if __name__ == "__main__":
    exit(main())